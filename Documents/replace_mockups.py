import sys
from pptx import Presentation
from pptx.util import Inches
import os

def replace_images_in_pptx(pptx_path, screenshots_dir, output_path):
    print(f"Opening {pptx_path}...")
    prs = Presentation(pptx_path)
    
    slide_map = {
        8: "screen1.png",
        9: "screen2.png",
        11: "screen3.png",
        12: "screen4.png",
        13: "screen5.png",
        14: "screen6.png",
        15: "screen7.png",
        16: "screen8.png",
        17: "screen9.png",
        18: "screen10.png",
        19: "screen11.png"
    }

    for slide_idx, filename in slide_map.items():
        if slide_idx < len(prs.slides):
            slide = prs.slides[slide_idx]
            img_path = os.path.join(screenshots_dir, filename)
            
            if not os.path.exists(img_path):
                print(f"Warning: {img_path} not found. Skipping slide {slide_idx+1}.")
                continue
                
            # The mockup container is usually around top=1.5, left=0.6, width=12.1, height=5.5
            container_left = 0.6
            container_top = 1.5
            container_right = 0.6 + 12.1
            container_bottom = 1.5 + 5.5
            
            # Find shapes that fall within or overlap significantly with this container
            shapes_to_delete = []
            for shape in slide.shapes:
                # Skip headers (top < 1.4)
                if shape.top.inches < 1.4:
                    continue
                # Skip slide numbers / footers (usually small text boxes at the very bottom or not fully inside)
                # But actually, the mockup has text boxes. Let's delete shapes whose center is inside the container
                center_x = shape.left.inches + (shape.width.inches / 2)
                center_y = shape.top.inches + (shape.height.inches / 2)
                
                if (container_left - 0.2 <= center_x <= container_right + 0.2) and \
                   (container_top - 0.2 <= center_y <= container_bottom + 0.2):
                    # Exclude slide number. Slide number is usually just a number, let's check text
                    is_slide_num = False
                    if shape.has_text_frame and shape.text.strip().isdigit():
                        is_slide_num = True
                    if not is_slide_num:
                        shapes_to_delete.append(shape)
                        
            print(f"Slide {slide_idx+1}: deleting {len(shapes_to_delete)} mockup shapes...")
            for shape in shapes_to_delete:
                sp = shape._element
                sp.getparent().remove(sp)
                
            # Insert the screenshot
            # To fit within 12.1 x 5.5 without distortion, we scale it
            # Screenshot is 16:9, meaning height = width * (9/16)
            # Max width = 12.1, height would be 12.1 * 9/16 = 6.8 (exceeds 5.5)
            # So max height = 5.5, width would be 5.5 * 16/9 = 9.77
            target_height = Inches(5.5)
            target_width = target_height * (16/9)
            
            target_top = Inches(1.5)
            target_left = Inches(0.6) + (Inches(12.1) - target_width) / 2 # Center horizontally
            
            print(f"Inserting {filename}...")
            slide.shapes.add_picture(img_path, target_left, target_top, target_width, target_height)

    print(f"Saving updated presentation to {output_path}...")
    prs.save(output_path)
    print("Done!")

if __name__ == "__main__":
    pptx_path = "SnapFlect_Product_Demo_Deck.pptx"
    screenshots_dir = "../snapflect-frontend/screenshots"
    output_path = "SnapFlect_Product_Demo_Deck_Updated.pptx"
    replace_images_in_pptx(pptx_path, screenshots_dir, output_path)
