import sys
from pptx import Presentation
from pptx.util import Inches

def debug_slide_pos(pptx_path, slide_idx):
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_idx]
    for i, shape in enumerate(slide.shapes):
        print(f"Shape {i}: name={shape.name}, top={shape.top.inches:.2f}, left={shape.left.inches:.2f}, width={shape.width.inches:.2f}, height={shape.height.inches:.2f}")
        if shape.has_text_frame:
            print(f"  Text: {shape.text[:50].replace(chr(10), ' ')}")

if __name__ == "__main__":
    debug_slide_pos("SnapFlect_Product_Demo_Deck.pptx", 9)
